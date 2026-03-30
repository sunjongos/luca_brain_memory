import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Title layout (usually index 0)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "남양주백병원 AI비서 NDB 활용 교육"
    subtitle.text = "Antigravity (솔로프리너 AI 비서) 100% 활용법\n자동화로 퇴근 시간을 앞당기자!"

    # Content layout (usually index 1)
    bullet_slide_layout = prs.slide_layouts[1]

    slides_data = [
        ("0단계: AI 비서(Antigravity) 다운로드 및 설치", [
            "컴퓨터를 잘 몰라도 괜찮습니다. 딱 3번만 클릭하세요!",
            "1. 검색: 구글(Google)에서 'Antigravity AI' 검색",
            "2. 다운로드: 공식 홈페이지 접속 후 'Windows 용 다운로드' 클릭",
            "3. 설치: 다운받은 설치 프로그램을 더블클릭 후 '다음(Next)'만 계속 누르기!",
            "설치가 끝나면 바탕화면에 우주인(?) 모양의 아이콘이 생깁니다."
        ]),
        ("1단계: 준비 작업 (권한 설정)", [
            "AI가 내 PC에서 원활하게 작업할 수 있도록 권한 부여",
            "Antigravity 실행: 아이콘 마우스 우클릭 -> '관리자 권한으로 실행'",
            "PowerShell 실행: 시작 메뉴에서 PowerShell 검색 후 '관리자 권한으로 실행'으로 한 개 열어두기",
            "이 작업만 해두면 AI가 프로그램을 설치할 때 오류 없이 한 번에 통과!"
        ]),
        ("2단계: NDB 비서 페르소나 부여", [
            "설정(Settings) 메뉴의 'Global Instruction' 칸에 아래 텍스트 복사/붙여넣기",
            "[이름] 남양주백병원 AI비서 NDB",
            "[역할] 업무(일정, 문서, 리서치)를 실질적으로 돕는 수석 비서",
            "[지침 1] 수동적 대답이 아닌 능동적/주도적 제안 (예: 진행할까요?)",
            "[지침 2] 의료 행정 보완 철저 유지 및 친절한 안내",
            "※ 코드 작성 없이 텍스트 복붙만으로 강력한 맞춤형 AI 비서 세팅 끝!"
        ]),
        ("3단계: 외부 두뇌 확장 (MCP 환경 세팅)", [
            "MCP(Model Context Protocol) 확장을 위한 필수 3대장: Node.js, Git, Python",
            "[AI 프롬프트 실습]",
            "\"내 PC에 MCP를 쓰기 위한 Node.js, Git, Python이 최신 버전으로 깔려있는지 확인하고, 없다면 알아서 설치해 줘.\"",
            "설치 후 NotebookLM MCP를 연결하여 병원 매뉴얼이나 방대한 자료 학습 및 질의응답 시연"
        ]),
        ("4단계: 내 업무와 연결 (Google Workspace 연동)", [
            "내 구글 일정, 문서, 이메일을 다룰 수 있는 권한(Credentials.json) 발급",
            "1. Google Cloud Console 접속 후 새 프로젝트 생성",
            "2. 5대 핵심 API 활성화: Gmail, Calendar, Docs, Slides, Drive",
            "3. OAuth 동의 화면에서 외부 사용자 및 테스트 이메일(본인 계정) 등록",
            "4. OAuth 클라이언트 ID 생성(데스크톱 앱) -> credentials.json 다운로드",
            "다운받은 파일을 바탕화면에 저장 후 AI에게 GWS CLI 연동 지시!"
        ]),
        ("5단계: 업무 자동화의 꽃 (Workflow & Skill)", [
            "매번 같은 지시를 반복하지 않고, 루틴한 업무를 자동화 시스템으로 구축하기",
            "[Workflow (워크플로우)]",
            "특정 작업 순서(아침 브리핑 등)를 .md 파일로 정리해 '이 순서대로 일해'라고 가르치는 방법",
            "[Skill (스킬)]",
            "웹 검색, 텔레그램 알림, 데이터베이스 조회 등 AI가 사용할 수 있는 '특수 능력' 장착",
            "프롬프트: \"방금 한 이메일 요약 작업을 /email_summary 워크플로우로 만들어줘.\""
        ])
    ]

    for title_text, bullets in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet

    v_end_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(v_end_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "감사합니다!"
    subtitle.text = "이제 여러분의 데스크탑에는 최고 수준의 비서가 상주합니다."

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "NDB_AI_비서_활용교육.pptx")
    prs.save(output_path)
    print(f"PPTX saved successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
